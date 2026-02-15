from __future__ import annotations

from datetime import datetime
from pathlib import Path
import textwrap

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parents[1]
DATA_PATH = ROOT / "data" / "hr_metrics_monthly.csv"
OUTPUT_DIR = ROOT / "output"


def load_data(path: Path) -> pd.DataFrame:
    if not path.exists():
        raise FileNotFoundError(f"Metrics file not found: {path}")
    df = pd.read_csv(path)
    df["month_date"] = pd.to_datetime(df["month"] + "-01")
    return df


def build_insights(df: pd.DataFrame) -> tuple[str, dict[str, pd.DataFrame]]:
    monthly = (
        df.groupby("month_date", as_index=False)
        .agg(
            headcount=("headcount", "sum"),
            new_hires=("new_hires", "sum"),
            terminations=("terminations", "sum"),
            open_positions=("open_positions", "sum"),
            avg_time_to_fill=("time_to_fill_days", "mean"),
            avg_offer_acceptance=("offer_acceptance_rate", "mean"),
            avg_turnover_rate=("turnover_rate", "mean"),
        )
        .sort_values("month_date")
    )

    monthly["headcount_mom_pct"] = monthly["headcount"].pct_change() * 100
    monthly["new_hires_mom_pct"] = monthly["new_hires"].pct_change() * 100
    monthly["turnover_mom_delta"] = monthly["avg_turnover_rate"].diff()

    latest = monthly.iloc[-1]
    prior = monthly.iloc[-2]

    dept_latest = (
        df[df["month_date"] == latest["month_date"]]
        .groupby("department", as_index=False)
        .agg(
            headcount=("headcount", "sum"),
            turnover_rate=("turnover_rate", "mean"),
            time_to_fill_days=("time_to_fill_days", "mean"),
        )
        .sort_values("headcount", ascending=False)
    )

    top_turnover = dept_latest.sort_values("turnover_rate", ascending=False).iloc[0]
    best_ttf = dept_latest.sort_values("time_to_fill_days", ascending=True).iloc[0]

    exec_summary = (
        f"As of {latest['month_date'].strftime('%b %Y')}, total headcount is {int(latest['headcount'])} "
        f"({latest['headcount_mom_pct']:+.1f}% MoM), with {int(latest['new_hires'])} hires and "
        f"{int(latest['terminations'])} terminations. Average turnover is {latest['avg_turnover_rate']:.2f}% "
        f"({latest['turnover_mom_delta']:+.2f} pts MoM), offer acceptance is {latest['avg_offer_acceptance']:.1f}%, "
        f"and time-to-fill averages {latest['avg_time_to_fill']:.1f} days. Highest turnover is in {top_turnover['department']} "
        f"({top_turnover['turnover_rate']:.2f}%), while fastest hiring cycle is in {best_ttf['department']} "
        f"({best_ttf['time_to_fill_days']:.1f} days)."
    )

    return exec_summary, {
        "monthly": monthly,
        "dept_latest": dept_latest,
        "latest": latest.to_frame().T,
        "prior": prior.to_frame().T,
    }


def save_chart_headcount(monthly: pd.DataFrame, out_path: Path) -> None:
    plt.figure(figsize=(10, 5))
    plt.plot(monthly["month_date"], monthly["headcount"], marker="o", linewidth=2.5)
    plt.title("Total Headcount Trend")
    plt.xlabel("Month")
    plt.ylabel("Employees")
    plt.grid(alpha=0.25)
    plt.tight_layout()
    plt.savefig(out_path, dpi=160)
    plt.close()


def save_chart_hiring(monthly: pd.DataFrame, out_path: Path) -> None:
    fig, ax1 = plt.subplots(figsize=(10, 5))
    ax1.plot(monthly["month_date"], monthly["new_hires"], marker="o", label="New Hires", color="#1f77b4")
    ax1.set_ylabel("New Hires", color="#1f77b4")
    ax1.tick_params(axis="y", labelcolor="#1f77b4")
    ax1.set_xlabel("Month")

    ax2 = ax1.twinx()
    ax2.plot(monthly["month_date"], monthly["avg_turnover_rate"], marker="s", label="Turnover Rate", color="#d62728")
    ax2.set_ylabel("Turnover Rate (%)", color="#d62728")
    ax2.tick_params(axis="y", labelcolor="#d62728")

    plt.title("Hiring Activity vs Turnover Rate")
    fig.tight_layout()
    plt.savefig(out_path, dpi=160)
    plt.close()


def save_chart_department(dept_latest: pd.DataFrame, out_path: Path) -> None:
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))

    axes[0].bar(dept_latest["department"], dept_latest["headcount"], color="#2ca02c")
    axes[0].set_title("Headcount by Department")
    axes[0].set_ylabel("Employees")
    axes[0].tick_params(axis="x", rotation=35)

    axes[1].bar(dept_latest["department"], dept_latest["turnover_rate"], color="#ff7f0e")
    axes[1].set_title("Turnover Rate by Department")
    axes[1].set_ylabel("Turnover Rate (%)")
    axes[1].tick_params(axis="x", rotation=35)

    plt.tight_layout()
    plt.savefig(out_path, dpi=160)
    plt.close()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.name = "Calibri"

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.5))
        stf = subtitle_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.name = "Calibri"


def build_presentation(exec_summary: str, frames: dict[str, pd.DataFrame], charts: dict[str, Path]) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title + executive summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "HR Metrics Automation Dashboard", datetime.now().strftime("Generated %b %d, %Y"))

    summary_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(4.8))
    tf = summary_box.text_frame
    tf.word_wrap = True

    header = tf.paragraphs[0]
    header.text = "Executive Summary"
    header.font.size = Pt(20)
    header.font.bold = True

    body = tf.add_paragraph()
    body.text = textwrap.fill(exec_summary, width=120)
    body.font.size = Pt(15)
    body.level = 0

    # Slide 2: Headcount trends
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide2, "Headcount Trends", "Monthly total headcount across all departments")
    slide2.shapes.add_picture(str(charts["headcount"]), Inches(0.7), Inches(1.6), width=Inches(12.0))

    # Slide 3: Hiring metrics
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide3, "Hiring Metrics", "New hires and turnover rate by month")
    slide3.shapes.add_picture(str(charts["hiring"]), Inches(0.7), Inches(1.6), width=Inches(12.0))

    # Slide 4: Department breakdown
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide4, "Departmental Breakdown", "Latest month headcount and turnover comparison")
    slide4.shapes.add_picture(str(charts["department"]), Inches(0.7), Inches(1.6), width=Inches(12.0))

    # Slide 5: Key insights + recommendations
    latest = frames["latest"].iloc[0]
    dept_latest = frames["dept_latest"]
    risk_dept = dept_latest.sort_values("turnover_rate", ascending=False).iloc[0]
    low_hc_dept = dept_latest.sort_values("headcount", ascending=True).iloc[0]

    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide5, "Key Insights & Recommendations")
    box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.5))
    tf5 = box.text_frame
    tf5.word_wrap = True

    bullets = [
        f"Total headcount reached {int(latest['headcount'])}, with momentum of {latest['headcount_mom_pct']:+.1f}% month-over-month.",
        f"Average turnover is {latest['avg_turnover_rate']:.2f}%; highest risk area is {risk_dept['department']} ({risk_dept['turnover_rate']:.2f}%).",
        f"Average time-to-fill is {latest['avg_time_to_fill']:.1f} days; focus pipeline acceleration for teams above 40 days.",
        f"Open positions currently total {int(latest['open_positions'])}; prioritize hiring plans for {low_hc_dept['department']} to avoid service bottlenecks.",
        "Recommendation: Launch targeted retention interventions for high-turnover departments and monitor offer acceptance weekly to reduce backfill lag.",
    ]

    first = tf5.paragraphs[0]
    first.text = bullets[0]
    first.font.size = Pt(16)

    for item in bullets[1:]:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)

    out_pptx = OUTPUT_DIR / "hr_metrics_executive_dashboard.pptx"
    prs.save(out_pptx)
    return out_pptx


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    df = load_data(DATA_PATH)
    exec_summary, frames = build_insights(df)
    monthly = frames["monthly"]
    dept_latest = frames["dept_latest"]

    headcount_chart = OUTPUT_DIR / "headcount_trend.png"
    hiring_chart = OUTPUT_DIR / "hiring_turnover_trend.png"
    dept_chart = OUTPUT_DIR / "department_breakdown.png"

    save_chart_headcount(monthly, headcount_chart)
    save_chart_hiring(monthly, hiring_chart)
    save_chart_department(dept_latest, dept_chart)

    pptx_path = build_presentation(
        exec_summary,
        frames,
        {
            "headcount": headcount_chart,
            "hiring": hiring_chart,
            "department": dept_chart,
        },
    )

    summary_path = OUTPUT_DIR / "executive_summary.txt"
    summary_path.write_text(exec_summary + "\n", encoding="utf-8")

    print("Executive Summary:")
    print(exec_summary)
    print(f"\nGenerated:\n- {pptx_path}\n- {summary_path}\n- {headcount_chart}\n- {hiring_chart}\n- {dept_chart}")


if __name__ == "__main__":
    main()
